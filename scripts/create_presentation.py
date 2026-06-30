"""
Présentation Projet 6 — Scoring Crédit « Prêt à Dépenser »
Style inspiré du template Projet 5 : fond blanc, titres bleus, diagrammes épurés.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import numpy as np
import io

from pathlib import Path

plt.rcParams["font.family"] = "DejaVu Sans"

# Racine du projet : scripts/ -> racine/
_ROOT    = Path(__file__).resolve().parent.parent

# Dossier des figures générées par le notebook
_FIG_DIR = _ROOT / "reports" / "figures"

def load_img(sl, filename, x, y, w, h):
    """
    Charge un PNG depuis reports/figures/ si disponible.
    Retourne True si le fichier existe, False sinon.
    """
    p = _FIG_DIR / filename
    if p.exists():
        sl.shapes.add_picture(str(p), x, y, w, h)
        return True
    return False

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x32, 0x64)   # bleu nuit (overlay titre)
BLUE    = RGBColor(0x2D, 0x6A, 0xB5)   # bleu principal (titres)
BLUE2   = RGBColor(0x5B, 0x9B, 0xD5)   # bleu clair (éléments)
TEAL    = RGBColor(0x17, 0xA5, 0xC8)   # accent cyan (ligne déco)
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # vert (succès)
ORANGE  = RGBColor(0xE6, 0x7E, 0x22)   # orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x66, 0x66, 0x66)
LGRAY   = RGBColor(0xF5, 0xF7, 0xFA)   # gris clair (fond secondaire)
TEXTE   = RGBColor(0x1A, 0x1A, 0x2E)

HN  = "#1F3264"   # navy
HB  = "#2D6AB5"   # blue
HB2 = "#5B9BD5"   # blue2
HT  = "#17A5C8"   # teal
HG  = "#27AE60"   # green
HO  = "#E67E22"   # orange
HW  = "#FFFFFF"
HLG = "#F5F7FA"

W = Inches(13.33)
H = Inches(7.5)

# ── Helpers pptx ──────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(sl, x, y, w, h, fill, border=None, bw=0.75, radius=False):
    sp = sl.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if border:
        sp.line.color.rgb = border
        sp.line.width = Pt(bw)
    else:
        sp.line.fill.background()
    return sp

def txt(sl, s, x, y, w, h, size=16, bold=False, color=TEXTE,
        align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = s
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def multiline(sl, lines, x, y, w, h, default_size=14, default_color=TEXTE,
              default_align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        text   = item[0]
        size   = item[1] if len(item) > 1 else default_size
        bold   = item[2] if len(item) > 2 else False
        color  = item[3] if len(item) > 3 else default_color
        align  = item[4] if len(item) > 4 else default_align
        space  = item[5] if len(item) > 5 else 2
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box

def fig_img(sl, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", transparent=True)
    buf.seek(0)
    sl.shapes.add_picture(buf, x, y, w, h)
    plt.close(fig)

def slide_header(sl, title, subtitle=None):
    """Titre bleu en haut à gauche avec ligne d'accent, sur fond blanc."""
    txt(sl, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.75),
        size=32, bold=True, color=BLUE)
    if subtitle:
        txt(sl, subtitle, Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
            size=14, color=GRAY, italic=True)
    rect(sl, Inches(0.5), Inches(1.05), Inches(3.5), Inches(0.05), TEAL)

def footer(sl, num):
    txt(sl, f"Faïza KHELLADI — Scoring Crédit | Prêt à Dépenser",
        Inches(0.5), H - Inches(0.38), Inches(10), Inches(0.35),
        size=9, color=GRAY)
    txt(sl, str(num), W - Inches(0.7), H - Inches(0.38), Inches(0.5), Inches(0.35),
        size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ── Figures matplotlib ────────────────────────────────────────────────────────

def fig_network_bg():
    """Fond réseau de nœuds pour la slide titre."""
    np.random.seed(7)
    fig, ax = plt.subplots(figsize=(13.33, 7.5), facecolor="#F0F2F5")
    n = 30
    pos = np.random.rand(n, 2) * np.array([13.33, 7.5])
    for i in range(n):
        for j in range(i+1, n):
            d = np.linalg.norm(pos[i] - pos[j])
            if d < 2.8:
                ax.plot([pos[i,0], pos[j,0]], [pos[i,1], pos[j,1]],
                        color="#B0BEC5", alpha=0.5, linewidth=0.8)
    sizes  = np.random.choice([30, 80, 180, 350, 600], n)
    clist  = (["#B0BEC5"] * 22 + ["#78909C"] * 4 + [HT] * 2 + [HB2] * 2)
    np.random.shuffle(clist)
    ax.scatter(pos[:,0], pos[:,1], s=sizes, c=clist, alpha=0.7, linewidths=0,
               zorder=3)
    ax.set_xlim(0, 13.33)
    ax.set_ylim(0, 7.5)
    ax.axis("off")
    fig.subplots_adjust(0, 0, 1, 1)
    return fig

def fig_pie_imbalance():
    fig, ax = plt.subplots(figsize=(4.5, 3.8), facecolor="none")
    vals   = [91.9, 8.1]
    labels = ["Non-défaut\n91.9 %", "Défaut\n8.1 %"]
    ax.pie(vals, labels=labels, colors=[HB2, HO], startangle=90,
           wedgeprops={"edgecolor": "white", "linewidth": 2.5},
           textprops={"fontsize": 12, "fontweight": "bold"})
    ax.set_title("Distribution des classes", fontsize=13,
                 fontweight="bold", pad=10, color=HN)
    return fig

def fig_contexte_diagram():
    """Diagramme sources → FE → Modèle → Score."""
    fig, ax = plt.subplots(figsize=(10, 3.5), facecolor="none")
    ax.axis("off")

    boxes = [
        (0.07, "7 tables\nsources", HN),
        (0.30, "Feature\nEngineering\n552 features", HB),
        (0.56, "Modèle ML\n(LightGBM)", HB2),
        (0.80, "Score\nde risque", HG),
    ]
    for x, label, color in boxes:
        bp = FancyBboxPatch((x - 0.1, 0.22), 0.185, 0.55,
                            boxstyle="round,pad=0.02",
                            facecolor=color, edgecolor="white", linewidth=1.5,
                            transform=ax.transAxes)
        ax.add_patch(bp)
        ax.text(x + 0.0, 0.49, label, ha="center", va="center",
                fontsize=11, color="white", fontweight="bold",
                transform=ax.transAxes)
        if x != 0.80:
            ax.annotate("", xy=(x + 0.125, 0.50), xytext=(x + 0.09, 0.50),
                        xycoords="axes fraction", textcoords="axes fraction",
                        arrowprops=dict(arrowstyle="->", color="#444", lw=1.8))

    sub = ["Bureau", "Credit card", "Prev. apps", "Installments", "POS Cash",
           "Application", "Balance"]
    for i, s in enumerate(sub):
        ax.text(0.015, 0.88 - i * 0.135, f"• {s}", fontsize=8, color="#555",
                transform=ax.transAxes)

    ax.text(0.40, 0.06, "Validation croisée · Seuil optimal · Coût métier",
            ha="center", fontsize=10, color=HB, fontstyle="italic",
            transform=ax.transAxes)
    fig.tight_layout()
    return fig

def fig_auc_step3():
    fig, ax = plt.subplots(figsize=(6, 3.8), facecolor="none")
    models = ["Dummy", "LogReg\nbaseline", "LightGBM\nbaseline"]
    aucs   = [0.500, 0.763, 0.780]
    colors = ["#B0BEC5", HB2, HN]
    bars   = ax.bar(models, aucs, color=colors, width=0.45,
                    edgecolor="white", linewidth=1.5)
    ax.axhline(0.5, linestyle="--", color=HO, linewidth=1.5, label="Aléatoire")
    ax.set_ylim(0, 0.95)
    ax.set_ylabel("AUC (OOF)", fontsize=11)
    ax.set_title("AUC — Étape 3 (baseline)", fontsize=13,
                 fontweight="bold", color=HN)
    for bar, auc in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width() / 2, auc + 0.012,
                f"{auc:.3f}", ha="center", fontsize=11,
                fontweight="bold", color=HN)
    ax.legend(fontsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("none")
    fig.tight_layout()
    return fig

def fig_business_cost():
    fig, axes = plt.subplots(1, 2, figsize=(8, 3.2), facecolor="none")
    ax1, ax2 = axes
    ax1.text(0.5, 0.65, "10 × FN + 1 × FP", ha="center", va="center",
             fontsize=26, fontweight="bold", color=HN,
             transform=ax1.transAxes)
    ax1.text(0.5, 0.35, "Minimisation du coût métier", ha="center",
             fontsize=13, color="#555", transform=ax1.transAxes)
    ax1.text(0.5, 0.12,
             "FN = défaut non détecté   |   FP = bon client refusé",
             ha="center", fontsize=9, color="#888",
             transform=ax1.transAxes)
    ax1.axis("off"); ax1.set_facecolor("none")
    ax2.bar(["Faux Négatif\n(FN)", "Faux Positif\n(FP)"], [10, 1],
            color=["#C0392B", HO], edgecolor="white", linewidth=1.5, width=0.4)
    ax2.set_title("Coût unitaire relatif", fontsize=12,
                  fontweight="bold", color=HN)
    ax2.set_ylabel("Coût"); ax2.set_ylim(0, 13)
    ax2.text(0, 10.5, "×10", ha="center", fontsize=13,
             fontweight="bold", color="#C0392B")
    ax2.text(1, 1.5, "×1", ha="center", fontsize=13,
             fontweight="bold", color=HO)
    ax2.spines[["top", "right"]].set_visible(False)
    ax2.set_facecolor("none")
    fig.tight_layout()
    return fig

def fig_auc_compare():
    fig, ax = plt.subplots(figsize=(7, 3.8), facecolor="none")
    labels  = ["LR\nbaseline", "LR\noptimisée", "LGB\nbaseline", "LGB\noptimisé"]
    aucs    = [0.763, 0.765, 0.780, 0.785]
    colors  = [HB2, HG, HN, HG]
    bars    = ax.bar(labels, aucs, color=colors, width=0.45,
                     edgecolor="white", linewidth=1.5)
    ax.set_ylim(0.72, 0.81)
    ax.set_ylabel("AUC (OOF)", fontsize=11)
    ax.set_title("Baseline vs Optimisé", fontsize=13,
                 fontweight="bold", color=HN)
    for bar, auc in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width() / 2, auc + 0.0003,
                f"{auc:.3f}", ha="center", fontsize=10,
                fontweight="bold", color=HN)
    p1 = mpatches.Patch(color=HB2, label="Baseline")
    p2 = mpatches.Patch(color=HG, label="Optimisé")
    ax.legend(handles=[p1, p2], fontsize=10)
    ax.spines[["top", "right"]].set_visible(False)
    ax.set_facecolor("none")
    fig.tight_layout()
    return fig

def fig_mlflow_pipeline():
    fig, ax = plt.subplots(figsize=(11, 2.8), facecolor="none")
    ax.axis("off"); ax.set_facecolor("none")
    steps = [
        ("Données\ntraitées", HN),
        ("CV\n5 folds", HB),
        ("Metriques\nOOF", HB2),
        ("MLflow\nRun", HG),
        ("UI MLflow\nvisualisation", HT),
    ]
    xs = [0.08, 0.27, 0.50, 0.72, 0.91]
    for (label, color), x in zip(steps, xs):
        bp = FancyBboxPatch((x - 0.09, 0.15), 0.16, 0.70,
                            boxstyle="round,pad=0.02",
                            facecolor=color, edgecolor="white", linewidth=2,
                            transform=ax.transAxes)
        ax.add_patch(bp)
        ax.text(x, 0.52, label, ha="center", va="center", fontsize=11,
                fontweight="bold", color="white", transform=ax.transAxes)
        if x != xs[-1]:
            ax.annotate("", xy=(x + 0.11, 0.50), xytext=(x + 0.09, 0.50),
                        xycoords="axes fraction", textcoords="axes fraction",
                        arrowprops=dict(arrowstyle="->", color="#444", lw=2))
    logged = ["params", "metriques", "figures (.png)", "modeles (.pkl + MLflow)"]
    for i, item in enumerate(logged):
        ax.text(0.47 + i * 0.135, 0.05, f"+ {item}", ha="center",
                fontsize=8, color=HG, transform=ax.transAxes, fontweight="bold")
    fig.tight_layout()
    return fig

def fig_optim_pipeline():
    fig, ax = plt.subplots(figsize=(11, 3.0), facecolor="none")
    ax.axis("off"); ax.set_facecolor("none")
    steps = [
        ("GridSearchCV\nLogReg", HB2, "C: 0.001/0.01/0.1\npenalty: l2\n20% dataset"),
        ("RandomizedSearch\nLightGBM", HB, "30 trials\nnum_leaves, lr\nsubsample..."),
        ("GridSearchCV\nfine-tuning", HN, "Exploration locale\nautour meilleur\nresultat"),
        ("CV Finale\n5 folds", HG, "Early stopping\nPredictions OOF\nMetriques + MLflow"),
    ]
    xs = [0.11, 0.36, 0.62, 0.87]
    for (label, color, detail), x in zip(steps, xs):
        bp = FancyBboxPatch((x - 0.115, 0.52), 0.20, 0.42,
                            boxstyle="round,pad=0.02",
                            facecolor=color, edgecolor="white", linewidth=1.5,
                            transform=ax.transAxes)
        ax.add_patch(bp)
        ax.text(x, 0.74, label, ha="center", va="center", fontsize=10,
                fontweight="bold", color="white", transform=ax.transAxes)
        ax.text(x, 0.30, detail, ha="center", va="center", fontsize=8.5,
                color="#555", transform=ax.transAxes)
        if x != xs[-1]:
            ax.annotate("", xy=(x + 0.14, 0.73), xytext=(x + 0.09, 0.73),
                        xycoords="axes fraction", textcoords="axes fraction",
                        arrowprops=dict(arrowstyle="->", color="#444", lw=2))
    fig.tight_layout()
    return fig

# ── Slides ─────────────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank(prs)
    # Fond réseau
    fig = fig_network_bg()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
    buf.seek(0)
    sl.shapes.add_picture(buf, 0, 0, W, H)
    plt.close(fig)

    # Overlay sombre (droite)
    rect(sl, Inches(4.8), Inches(2.2), Inches(8.2), Inches(3.5), NAVY)

    # Textes
    txt(sl, "Formation Data Scientist Machine Learning",
        Inches(5.1), Inches(2.4), Inches(7.7), Inches(0.45),
        size=13, color=RGBColor(0xB0, 0xC8, 0xE8))
    txt(sl, "Projet 6 :", Inches(5.1), Inches(2.82), Inches(7.7), Inches(0.45),
        size=15, bold=True, color=RGBColor(0xB0, 0xC8, 0xE8))
    txt(sl, "Scoring Credit — Pret a Depenser",
        Inches(5.1), Inches(3.2), Inches(7.7), Inches(0.9),
        size=28, bold=True, color=WHITE)
    # Ligne accent cyan
    rect(sl, Inches(5.1), Inches(4.15), Inches(7.2), Inches(0.05), TEAL)
    txt(sl, "Faiza KHELLADI", Inches(5.1), Inches(4.35), Inches(7.5), Inches(0.5),
        size=17, color=WHITE)
    return sl

def slide_sommaire(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Sommaire")
    footer(sl, 2)

    items = [
        "Contexte & Problematique",
        "Donnees & Feature Engineering",
        "Metrique Metier",
        "Pipeline de Modelisation",
        "Etape 3 — Comparaison des Modeles Baseline",
        "Etape 4 — Optimisation des Hyperparametres",
        "Tracking MLflow",
        "Conclusion & Perspectives",
    ]
    y = Inches(1.5)
    for item in items:
        rect(sl, Inches(0.9), y + Inches(0.12), Inches(0.12), Inches(0.12), BLUE)
        txt(sl, item, Inches(1.2), y, Inches(10), Inches(0.4),
            size=16, color=TEXTE)
        y += Inches(0.58)
    return sl

def slide_contexte(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Contexte & Problematique")
    footer(sl, 3)

    # ── Colonne gauche : description métier ────────────────────────────
    multiline(sl, [
        ("Societe : Pret a Depenser", 16, True, NAVY),
        ("", 5),
        ("Etablissement financier proposant des credits a la consommation",
         13, False, TEXTE),
        ("a des personnes peu ou pas bancarisees.", 13, False, TEXTE),
        ("", 10),
        ("Objectif du projet", 15, True, BLUE),
        ("", 5),
        ("Construire un modele de scoring credit pour predire la", 13, False, TEXTE),
        ("probabilite de defaut de remboursement de chaque client,", 13, False, TEXTE),
        ("et definir un seuil de decision optimal.", 13, False, TEXTE),
        ("", 10),
        ("Problematique metier", 15, True, BLUE),
        ("", 5),
        ("Un client defaillant non detecte (faux negatif) coute", 13, False, TEXTE),
        ("10 fois plus qu'un bon client refuse (faux positif).", 13, False, TEXTE),
        ("Le modele doit minimiser ce cout asymetrique.", 13, False, TEXTE),
        ("", 10),
        ("Donnees : 307 511 clients  |  7 tables sources  |  552 features",
         12, False, GRAY),
    ], Inches(0.5), Inches(1.35), Inches(7.8), Inches(5.9))

    # ── Colonne droite : chiffres clés ─────────────────────────────────
    # Grand chiffre 1
    rect(sl, Inches(8.7), Inches(1.35), Inches(4.3), Inches(1.55), NAVY)
    txt(sl, "307 511", Inches(8.7), Inches(1.42), Inches(4.3), Inches(0.9),
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "clients dans le dataset", Inches(8.7), Inches(2.23), Inches(4.3), Inches(0.45),
        size=12, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.CENTER)

    # Grand chiffre 2
    rect(sl, Inches(8.7), Inches(3.05), Inches(4.3), Inches(1.55), ORANGE)
    txt(sl, "8.1 %", Inches(8.7), Inches(3.12), Inches(4.3), Inches(0.9),
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "taux de defaut (classes desequilibrees)", Inches(8.7), Inches(3.93),
        Inches(4.3), Inches(0.45), size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Pie chart
    fig = fig_pie_imbalance()
    fig_img(sl, fig, Inches(8.7), Inches(4.7), Inches(4.3), Inches(2.5))

    return sl

def slide_donnees(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Donnees & Feature Engineering")
    footer(sl, 4)

    # ── 7 tables : 2 lignes (4 + 3), zone gauche ──────────────────────
    # Chaque carte : nom de table (bold) + 2-3 colonnes clés
    tables = [
        ("Application\nprincipale",  "AMT_CREDIT\nAMT_INCOME_TOTAL\nDAYS_BIRTH",        NAVY),
        ("Bureau\ncredit",           "CREDIT_TYPE\nCREDIT_ACTIVE\nAMT_CREDIT_SUM",       BLUE),
        ("Historique\nbureau",       "MONTHS_BALANCE\nSTATUS (retard DPD)",              BLUE2),
        ("Applications\nprecedentes","AMT_APPLICATION\nNAME_CONTRACT_STATUS\nAMT_ANNUITY", NAVY),
        ("Soldes\nPOS / Cash",       "SK_DPD\nCNT_INSTALMENT\nMONTHS_BALANCE",           BLUE),
        ("Paiements\nechelonnes",    "AMT_INSTALMENT\nAMT_PAYMENT\nDAYS_ENTRY_PAYMENT",  BLUE2),
        ("Carte\ncredit",            "AMT_BALANCE\nAMT_DRAWINGS_CURRENT\nAMT_PAYMENT_CURRENT", NAVY),
    ]

    CW   = Inches(1.9)   # largeur carte
    CH   = Inches(2.05)  # hauteur carte
    STEP = Inches(2.15)  # pas horizontal (carte + gap)

    # Ligne 1 : 4 cartes
    y1 = Inches(1.35)
    for i, (name, cols, color) in enumerate(tables[:4]):
        x = Inches(0.4) + i * STEP
        rect(sl, x, y1, CW, CH, color)
        txt(sl, name, x, y1 + Inches(0.1), CW, Inches(0.65),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, cols,  x, y1 + Inches(0.78), CW, Inches(1.15),
            size=8.5, color=RGBColor(0xD0, 0xE8, 0xFF), align=PP_ALIGN.CENTER)

    # Ligne 2 : 3 cartes centrées sous la ligne 1
    y2  = y1 + CH + Inches(0.18)
    # Largeur totale ligne 1 : 4 * STEP - gap = 4*2.15 - 0.25 = 8.35"
    # Centre ligne 2 : 3 * STEP - gap = 6.2", offset = (8.35 - 6.2)/2 = 1.075"
    x2_start = Inches(0.4) + Inches(1.075)
    for i, (name, cols, color) in enumerate(tables[4:]):
        x = x2_start + i * STEP
        rect(sl, x, y2, CW, CH, color)
        txt(sl, name, x, y2 + Inches(0.1), CW, Inches(0.65),
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, cols,  x, y2 + Inches(0.78), CW, Inches(1.15),
            size=8.5, color=RGBColor(0xD0, 0xE8, 0xFF), align=PP_ALIGN.CENTER)

    # ── Bande Feature Engineering en bas ──────────────────────────────
    y_fe = y2 + CH + Inches(0.15)
    rect(sl, Inches(0.4), y_fe, Inches(8.6), Inches(0.85), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))
    multiline(sl, [
        ("Feature Engineering  —  552 features", 11, True, BLUE),
        ("Agregations stat. (mean/max/min/std) · Ratios financiers · "
         "Retards de paiement · Encodage categoriel", 10, False, TEXTE),
    ], Inches(0.6), y_fe + Inches(0.06), Inches(8.2), Inches(0.75))

    # ── Colonne droite : pie chart + note ─────────────────────────────
    fig = fig_pie_imbalance()
    fig_img(sl, fig, Inches(9.25), Inches(1.35), Inches(3.8), Inches(3.8))

    rect(sl, Inches(9.25), Inches(5.2), Inches(3.8), Inches(1.6), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))
    multiline(sl, [
        ("Desequilibre de classes", 12, True, BLUE),
        ("8.1 % de clients defaillants.", 11),
        ("class_weight='balanced' applique", 11),
        ("sur tous les modeles.", 11),
    ], Inches(9.45), Inches(5.3), Inches(3.5), Inches(1.4))

    return sl

def slide_pipeline(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Pipeline de Modelisation")
    footer(sl, 5)

    # Image générée depuis le diagramme utilisateur
    if not load_img(sl, 'pipeline_modelisation.png',
                    Inches(0.3), Inches(1.15), Inches(12.7), Inches(6.1)):
        # Fallback minimal si le fichier est absent
        txt(sl, "pipeline_modelisation.png introuvable",
            Inches(0.5), Inches(3.5), Inches(12), Inches(1.0),
            size=14, color=GRAY, align=PP_ALIGN.CENTER)
    return sl


def slide_metrique(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Metrique Metier & Evaluation", "Seuil optimal — Predictions OOF")
    footer(sl, 6)

    # Figure formule + barres
    fig = fig_business_cost()
    fig_img(sl, fig, Inches(0.5), Inches(1.4), Inches(7.5), Inches(3.4))

    # Pipeline evaluation
    rect(sl, Inches(0.5), Inches(5.1), Inches(12.5), Inches(1.9), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))
    multiline(sl, [
        ("Pipeline d'evaluation (OOF)  :  "
         "StratifiedKFold (5 folds)  →  Probabilites OOF  →  Seuil optimal  →  Metriques finales", 13, True, BLUE),
        ("", 4),
        ("Metriques suivies : AUC · PR-AUC · Recall · Precision · F1 · Seuil · Cout metier · FN · FP",
         12, False, TEXTE),
    ], Inches(0.8), Inches(5.2), Inches(12.1), Inches(1.6))

    # Boites explication
    items = [
        ("Pas de data leakage", "Fit sur le fold train\nseulement", BLUE),
        ("Seuil != 0.5", "Optimise sur les OOF\npour minimiser le cout", BLUE2),
        ("AUC pour GridSearch", "Independant du seuil\npour la recherche", NAVY),
    ]
    x = Inches(8.1)
    for title, desc, color in items:
        rect(sl, x, Inches(1.55), Inches(1.65), Inches(3.0), color)
        txt(sl, title, x + Inches(0.08), Inches(1.65), Inches(1.5), Inches(0.7),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, x + Inches(0.08), Inches(2.35), Inches(1.5), Inches(1.5),
            size=10, color=RGBColor(0xE0, 0xEA, 0xFF), align=PP_ALIGN.CENTER)
        x += Inches(1.75)
    return sl

def slide_step3(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Etape 3 — Comparaison des Modeles Baseline")
    footer(sl, 7)

    # Courbes ROC étape 3 — PNG réel si disponible, sinon bar chart placeholder
    if not load_img(sl, 'roc_step3.png', Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.6)):
        fig = fig_auc_step3()
        fig_img(sl, fig, Inches(0.5), Inches(1.3), Inches(6.5), Inches(5.6))

    # Tableau métriques
    rect(sl, Inches(7.2), Inches(1.3), Inches(5.9), Inches(5.6), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))

    # En-tête
    cols  = ["Modele", "AUC", "Recall", "Cout metier"]
    x_pos = [Inches(7.3), Inches(9.0), Inches(10.2), Inches(11.4)]
    col_w = [Inches(1.6), Inches(1.2), Inches(1.2), Inches(1.6)]
    rect(sl, Inches(7.2), Inches(1.3), Inches(5.9), Inches(0.55), NAVY)
    for col, xp, cw in zip(cols, x_pos, col_w):
        txt(sl, col, xp, Inches(1.33), cw, Inches(0.5),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Dummy",      "0.500", "0.089", "—",       WHITE),
        ("LogReg",     "0.763", "0.682", "124 580", WHITE),
        ("LightGBM",   "0.780", "0.695", "117 340",
         RGBColor(0xE8, 0xF5, 0xEE)),
    ]
    y = Inches(1.88)
    for model, auc, rec, cost, bg in rows:
        rect(sl, Inches(7.2), y, Inches(5.9), Inches(0.62), bg)
        for val, xp, cw in zip([model, auc, rec, cost], x_pos, col_w):
            txt(sl, val, xp, y + Inches(0.1), cw, Inches(0.45),
                size=12, color=TEXTE, align=PP_ALIGN.CENTER)
        y += Inches(0.64)

    # Encart LightGBM
    rect(sl, Inches(7.2), Inches(3.88), Inches(5.9), Inches(1.4), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))
    multiline(sl, [
        ("Pourquoi LightGBM ?", 13, True, BLUE),
        ("Gestion native des valeurs manquantes", 12),
        ("Acceleration GPU (device=gpu, OpenCL)", 12),
        ("Early stopping automatique (100 rounds)", 12),
    ], Inches(7.4), Inches(3.98), Inches(5.5), Inches(1.2))

    multiline(sl, [
        ("LightGBM domine des le baseline", 14, True, BLUE),
        ("+1.7 pts AUC vs LogReg   |   Cout metier reduit de 6 %", 12, False, GRAY),
    ], Inches(7.2), Inches(5.4), Inches(5.9), Inches(0.9))
    return sl

def slide_step4(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Etape 4 — Optimisation des Hyperparametres")
    footer(sl, 8)

    # Pipeline optimisation
    fig = fig_optim_pipeline()
    fig_img(sl, fig, Inches(0.5), Inches(1.25), Inches(12.5), Inches(2.8))

    # Courbe coût métier vs seuil — PNG réel si disponible, sinon bar chart placeholder
    if not load_img(sl, 'threshold_lgbm_opt.png', Inches(0.5), Inches(4.2), Inches(7.0), Inches(3.0)):
        fig2 = fig_auc_compare()
        fig_img(sl, fig2, Inches(0.5), Inches(4.2), Inches(7.0), Inches(3.0))

    # Résultats
    rect(sl, Inches(7.8), Inches(4.2), Inches(5.3), Inches(3.0), LGRAY,
         border=RGBColor(0xD0, 0xD8, 0xE8))
    multiline(sl, [
        ("Resultats de l'optimisation", 15, True, BLUE),
        ("", 5),
        ("LightGBM optimise : AUC ~ 0.785", 13, True, GREEN),
        ("Gain : +0.5 pts vs baseline LGB", 12),
        ("", 5),
        ("LR optimisee : AUC ~ 0.765", 13, True, BLUE),
        ("Gain marginal vs baseline LR", 12),
        ("", 5),
        ("Seuil optimal != 0.5 (OOF)", 12, False, GRAY),
        ("Cout metier minimise apres optimisation", 12, False, GRAY),
        ("Modeles sauvegardes en .pkl + MLflow", 12, False, GRAY),
    ], Inches(8.0), Inches(4.35), Inches(5.0), Inches(2.7))
    return sl

def slide_mlflow(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Tracking MLflow — Tracebilite des Experimentations")
    footer(sl, 9)

    # Pipeline
    fig = fig_mlflow_pipeline()
    fig_img(sl, fig, Inches(0.5), Inches(1.25), Inches(12.5), Inches(2.5))

    # 3 colonnes
    col_data = [
        ("Runs MLflow", NAVY, [
            "Dummy_baseline",
            "LR_baseline",
            "LightGBM_baseline",
            "LR_optimise",
            "LightGBM_optimise",
            "",
            "Tags : etape · model_type",
        ]),
        ("Metriques loggees", BLUE, [
            "auc_mean · auc_std",
            "pr_auc · recall",
            "precision · f1",
            "seuil · cout_metier",
            "fn · fp",
        ]),
        ("Artefacts sauvegardes", BLUE2, [
            "mlflow.sklearn.log_model (LR)",
            "mlflow.lightgbm.log_model (LGB)",
            "threshold_curve.png",
            "roc_curve.png",
            "Modeles .pkl locaux",
        ]),
    ]
    x = Inches(0.5)
    for title, color, items in col_data:
        rect(sl, x, Inches(4.0), Inches(4.1), Inches(3.15), color)
        txt(sl, title, x + Inches(0.15), Inches(4.07), Inches(3.8), Inches(0.55),
            size=14, bold=True, color=WHITE)
        lines = [(item, 12, False, WHITE) for item in items]
        multiline(sl, lines, x + Inches(0.2), Inches(4.65),
                  Inches(3.7), Inches(2.4))
        x += Inches(4.25)
    return sl

def slide_conclusion(prs):
    sl = blank(prs)
    rect(sl, 0, 0, W, H, WHITE)
    slide_header(sl, "Conclusion & Perspectives")
    footer(sl, 10)

    # Conclusion (gauche)
    rect(sl, Inches(0.5), Inches(1.35), Inches(6.0), Inches(0.52), NAVY)
    txt(sl, "Conclusion", Inches(0.5), Inches(1.38), Inches(6.0), Inches(0.46),
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.5), Inches(1.9), Inches(6.0), Inches(4.8), WHITE,
         border=RGBColor(0xB0, 0xC4, 0xD8))
    multiline(sl, [
        ("Meilleur modele : LightGBM optimise", 13, True, BLUE),
        ("AUC ~ 0.785 — superieur au baseline (+2.2 pts)", 12),
        ("", 5),
        ("Cout metier minimise", 13, True, BLUE),
        ("Seuil optimal calcule sur predictions OOF", 12),
        ("", 5),
        ("Pipeline reproductible", 13, True, BLUE),
        ("CV stratifiee · feature engineering centralise", 12),
        ("", 5),
        ("Tracking MLflow operationnel", 13, True, BLUE),
        ("5 runs loggues · modeles versiones · figures", 12),
        ("", 5),
        ("GPU active (LightGBM via OpenCL)", 13, True, BLUE),
        ("Acceleration de l'entrainement et de la recherche", 12),
    ], Inches(0.7), Inches(2.0), Inches(5.7), Inches(4.5))

    # Perspectives (droite)
    rect(sl, Inches(6.8), Inches(1.35), Inches(6.2), Inches(0.52), BLUE2)
    txt(sl, "Perspectives", Inches(6.8), Inches(1.38), Inches(6.2), Inches(0.46),
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, Inches(6.8), Inches(1.9), Inches(6.2), Inches(4.8), WHITE,
         border=RGBColor(0xB0, 0xC4, 0xD8))
    next_steps = [
        ("Feature Importance (SHAP)", "Interpretation globale et locale\ndes predictions clients"),
        ("API de scoring (FastAPI)", "Endpoint REST pour predictions\nen temps reel"),
        ("Dashboard client", "Visualisation du score et\ndes facteurs explicatifs"),
        ("MLOps continu", "Re-entrainement automatise\n+ monitoring de la derive"),
    ]
    y = Inches(2.05)
    for title, desc in next_steps:
        rect(sl, Inches(7.0), y + Inches(0.1), Inches(0.08), Inches(0.7), TEAL)
        txt(sl, title, Inches(7.25), y, Inches(5.5), Inches(0.42),
            size=13, bold=True, color=BLUE)
        txt(sl, desc, Inches(7.25), y + Inches(0.38), Inches(5.5), Inches(0.55),
            size=11, color=GRAY)
        y += Inches(1.08)
    return sl

# ── Main ───────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_title(prs)
    slide_sommaire(prs)
    slide_contexte(prs)
    slide_donnees(prs)
    slide_metrique(prs)
    slide_pipeline(prs)
    slide_step3(prs)
    slide_step4(prs)
    slide_mlflow(prs)
    slide_conclusion(prs)

    out = _ROOT / "reports" / "scoring_credit_presentation.pptx"
    prs.save(out)
    print(f"Presentation sauvegardee : {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
